"""
Document Tools - Tools for reading various document formats.
"""

from pathlib import Path
from typing import Optional

from .base import BaseTool, ToolResult


class ReadPDFTool(BaseTool):
    """Tool for reading PDF documents."""
    
    @property
    def name(self) -> str:
        return "read_pdf"
    
    @property
    def description(self) -> str:
        return "Extract text content from a PDF file."
    
    @property
    def parameters(self) -> dict:
        return {
            "type": "object",
            "properties": {
                "path": {
                    "type": "string",
                    "description": "Path to the PDF file"
                },
                "start_page": {
                    "type": "integer",
                    "description": "Start from this page (1-indexed, optional)"
                },
                "end_page": {
                    "type": "integer",
                    "description": "End at this page (inclusive, optional)"
                }
            },
            "required": ["path"]
        }
    
    def execute(self, **kwargs) -> ToolResult:
        path = kwargs.get("path")
        start_page = kwargs.get("start_page")
        end_page = kwargs.get("end_page")
        
        try:
            # Import PyPDF2 (optional dependency)
            try:
                from PyPDF2 import PdfReader
            except ImportError:
                return ToolResult.error_result(
                    "PyPDF2 is not installed. Run: pip install PyPDF2"
                )
            
            file_path = Path(path).expanduser().resolve()
            
            if not file_path.exists():
                return ToolResult.error_result(f"File not found: {path}")
            
            if file_path.suffix.lower() != ".pdf":
                return ToolResult.error_result(f"Not a PDF file: {path}")
            
            reader = PdfReader(str(file_path))
            total_pages = len(reader.pages)
            
            # Handle page range
            start_idx = (start_page - 1) if start_page else 0
            end_idx = end_page if end_page else total_pages
            
            # Clamp to valid range
            start_idx = max(0, min(start_idx, total_pages - 1))
            end_idx = max(start_idx + 1, min(end_idx, total_pages))
            
            # Extract text
            text_parts = []
            for i in range(start_idx, end_idx):
                page = reader.pages[i]
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(f"--- Page {i + 1} ---\n{page_text}")
            
            if not text_parts:
                return ToolResult.success_result(
                    f"PDF has {total_pages} pages but no extractable text.",
                    data={"pages": total_pages, "extracted": 0}
                )
            
            content = "\n\n".join(text_parts)
            return ToolResult.success_result(
                f"Extracted from {file_path} ({total_pages} pages total):\n\n{content}",
                data={"pages": total_pages, "extracted": end_idx - start_idx}
            )
            
        except Exception as e:
            return ToolResult.error_result(f"Error reading PDF: {str(e)}")


class ReadDocxTool(BaseTool):
    """Tool for reading Word documents."""
    
    @property
    def name(self) -> str:
        return "read_docx"
    
    @property
    def description(self) -> str:
        return "Extract text content from a Word (.docx) document."
    
    @property
    def parameters(self) -> dict:
        return {
            "type": "object",
            "properties": {
                "path": {
                    "type": "string",
                    "description": "Path to the Word document"
                },
                "include_tables": {
                    "type": "boolean",
                    "description": "Include table content",
                    "default": True
                }
            },
            "required": ["path"]
        }
    
    def execute(self, **kwargs) -> ToolResult:
        path = kwargs.get("path")
        include_tables = kwargs.get("include_tables", True)
        
        try:
            # Import python-docx (optional dependency)
            try:
                from docx import Document
            except ImportError:
                return ToolResult.error_result(
                    "python-docx is not installed. Run: pip install python-docx"
                )
            
            file_path = Path(path).expanduser().resolve()
            
            if not file_path.exists():
                return ToolResult.error_result(f"File not found: {path}")
            
            if file_path.suffix.lower() != ".docx":
                return ToolResult.error_result(f"Not a Word document: {path}")
            
            doc = Document(str(file_path))
            
            # Extract paragraphs
            text_parts = []
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)
            
            # Extract tables
            if include_tables and doc.tables:
                text_parts.append("\n--- Tables ---")
                for i, table in enumerate(doc.tables, 1):
                    table_text = [f"\nTable {i}:"]
                    for row in table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells)
                        table_text.append(row_text)
                    text_parts.append("\n".join(table_text))
            
            if not text_parts:
                return ToolResult.success_result(
                    "Document is empty or contains no extractable text.",
                    data={"paragraphs": 0, "tables": len(doc.tables)}
                )
            
            content = "\n\n".join(text_parts)
            return ToolResult.success_result(
                f"Extracted from {file_path}:\n\n{content}",
                data={
                    "paragraphs": len(doc.paragraphs),
                    "tables": len(doc.tables)
                }
            )
            
        except Exception as e:
            return ToolResult.error_result(f"Error reading Word document: {str(e)}")


class ReadPPTTool(BaseTool):
    """Tool for reading PowerPoint presentations."""

    @property
    def name(self) -> str:
        return "read_ppt"

    @property
    def description(self) -> str:
        return "Extract text content from a PowerPoint (.pptx) presentation."

    @property
    def parameters(self) -> dict:
        return {
            "type": "object",
            "properties": {
                "path": {
                    "type": "string",
                    "description": "Path to the PowerPoint file"
                },
                "include_notes": {
                    "type": "boolean",
                    "description": "Include speaker notes",
                    "default": True
                }
            },
            "required": ["path"]
        }

    def execute(self, **kwargs) -> ToolResult:
        path = kwargs.get("path")
        include_notes = kwargs.get("include_notes", True)

        try:
            # Import python-pptx (optional dependency)
            try:
                from pptx import Presentation
            except ImportError:
                return ToolResult.error_result(
                    "python-pptx is not installed. Run: pip install python-pptx"
                )

            file_path = Path(path).expanduser().resolve()

            if not file_path.exists():
                return ToolResult.error_result(f"File not found: {path}")

            if file_path.suffix.lower() not in (".pptx", ".ppt"):
                return ToolResult.error_result(f"Not a PowerPoint file: {path}")

            prs = Presentation(str(file_path))
            total_slides = len(prs.slides)

            # Extract text from slides
            text_parts = []
            for i, slide in enumerate(prs.slides, 1):
                slide_text = [f"--- Slide {i} ---"]

                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)

                # Extract speaker notes
                if include_notes and slide.has_notes_slide:
                    notes_frame = slide.notes_slide.notes_text_frame
                    if notes_frame and notes_frame.text.strip():
                        slide_text.append(f"\n[Speaker Notes]\n{notes_frame.text}")

                if len(slide_text) > 1:  # More than just the header
                    text_parts.append("\n".join(slide_text))

            if not text_parts:
                return ToolResult.success_result(
                    f"PowerPoint has {total_slides} slides but no extractable text.",
                    data={"slides": total_slides, "extracted": 0}
                )

            content = "\n\n".join(text_parts)
            return ToolResult.success_result(
                f"Extracted from {file_path} ({total_slides} slides):\n\n{content}",
                data={"slides": total_slides, "extracted": len(text_parts)}
            )

        except Exception as e:
            return ToolResult.error_result(f"Error reading PowerPoint: {str(e)}")


class ReadImageTool(BaseTool):
    """Tool for extracting text from images using OCR."""

    @property
    def name(self) -> str:
        return "read_image"
    
    @property
    def description(self) -> str:
        return "Extract text from an image using OCR (Optical Character Recognition)."
    
    @property
    def parameters(self) -> dict:
        return {
            "type": "object",
            "properties": {
                "path": {
                    "type": "string",
                    "description": "Path to the image file"
                },
                "language": {
                    "type": "string",
                    "description": "OCR language (e.g., 'eng', 'chi_sim')",
                    "default": "eng"
                }
            },
            "required": ["path"]
        }
    
    def execute(self, **kwargs) -> ToolResult:
        path = kwargs.get("path")
        language = kwargs.get("language", "eng")
        
        try:
            # Import dependencies
            try:
                from PIL import Image
                import pytesseract
            except ImportError as e:
                missing = str(e).split("'")[1]
                return ToolResult.error_result(
                    f"{missing} is not installed. Run: pip install Pillow pytesseract\n"
                    "Also ensure Tesseract OCR is installed on your system."
                )
            
            file_path = Path(path).expanduser().resolve()
            
            if not file_path.exists():
                return ToolResult.error_result(f"File not found: {path}")
            
            # Check if it's an image
            valid_extensions = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".webp"}
            if file_path.suffix.lower() not in valid_extensions:
                return ToolResult.error_result(
                    f"Not a supported image format: {path}. "
                    f"Supported: {', '.join(valid_extensions)}"
                )
            
            # Open and process image
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image, lang=language)
            
            if not text.strip():
                return ToolResult.success_result(
                    "No text detected in the image.",
                    data={"size": image.size, "text_found": False}
                )
            
            return ToolResult.success_result(
                f"Text extracted from {file_path}:\n\n{text}",
                data={
                    "size": image.size,
                    "text_found": True,
                    "char_count": len(text)
                }
            )
            
        except Exception as e:
            return ToolResult.error_result(f"Error reading image: {str(e)}")
